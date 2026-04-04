"""Extract metadata and to-dos from .pptx files."""

from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation


@dataclass
class PresentationMeta:
    title: str
    file_path: Path
    slide_count: int
    topics: list[str] = field(default_factory=list)
    todos: list[str] = field(default_factory=list)
    slide_summaries: list[tuple[str, str]] = field(default_factory=list)  # (title, body)


def extract(pptx_path: Path) -> PresentationMeta:
    """Extract metadata and to-dos from a PowerPoint file."""
    prs = Presentation(pptx_path)

    title = _extract_title(prs, pptx_path)
    topics = _extract_topics(prs)
    todos = _extract_todos(prs)
    slide_summaries = _extract_slide_summaries(prs)

    return PresentationMeta(
        title=title,
        file_path=pptx_path,
        slide_count=len(prs.slides),
        topics=topics,
        todos=todos,
        slide_summaries=slide_summaries,
    )


def _placeholder_idx(shape: object) -> int | None:
    """Return placeholder index or None if shape is not a placeholder."""
    try:
        fmt = shape.placeholder_format  # type: ignore[attr-defined]
        return fmt.idx if fmt is not None else None
    except Exception:
        return None


_GENERIC_TITLES = {"powerpoint presentation", "presentation", "untitled"}


def _first_slide_title(prs: Presentation) -> str:
    if prs.slides:
        for shape in prs.slides[0].shapes:
            if _placeholder_idx(shape) == 0 and shape.has_text_frame:
                return shape.text_frame.text.strip()
    return ""


def _extract_title(prs: Presentation, fallback: Path) -> str:
    """Use core properties title unless it's generic, then prefer first slide title."""
    core_title = (prs.core_properties.title or "").strip()
    if core_title and core_title.lower() not in _GENERIC_TITLES:
        return core_title

    slide_title = _first_slide_title(prs)
    if slide_title:
        return slide_title

    if core_title:
        return core_title

    return fallback.stem.replace("-", " ").replace("_", " ").title()


def _extract_topics(prs: Presentation) -> list[str]:
    """Extract slide titles as topic list."""
    topics: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if _placeholder_idx(shape) == 0 and shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    topics.append(text)
    return topics


def _extract_slide_summaries(prs: Presentation) -> list[tuple[str, str]]:
    """Extract (slide_title, body_text) for each slide."""
    summaries = []
    for slide in prs.slides:
        slide_title = ""
        body_parts: list[str] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            idx = _placeholder_idx(shape)
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if idx == 0:
                slide_title = text
            else:
                body_parts.append(text)

        body = " · ".join(body_parts)
        if slide_title or body:
            summaries.append((slide_title, body))

    return summaries


def _extract_todos(prs: Presentation) -> list[str]:
    """Extract lines containing TODO, Action, or follow-up keywords."""
    keywords = ("todo", "action item", "follow up", "follow-up", "next step", "to-do")
    todos: list[str] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if any(kw in text.lower() for kw in keywords):
                    todos.append(text)

    return todos
