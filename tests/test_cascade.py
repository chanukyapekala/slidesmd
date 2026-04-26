"""Tests for the cascading extraction logic in extractor.py."""

from __future__ import annotations

from pathlib import Path
from unittest.mock import patch, MagicMock

import pytest

from slidesmd.extractor import (
    PresentationMeta,
    _cascading_extract,
    _build_extractor_chain,
    extract,
)


def _make_meta(
    title: str = "Good Title",
    slide_count: int = 3,
    topics: list[str] | None = None,
    summaries: list[tuple[str, str]] | None = None,
    path: Path | None = None,
) -> PresentationMeta:
    """Helper to build PresentationMeta with defaults."""
    return PresentationMeta(
        title=title,
        file_path=path or Path("/tmp/deck.pptx"),
        slide_count=slide_count,
        topics=topics or ["A", "B", "C"],
        slide_summaries=summaries or [
            ("A", "Content A is fairly long to pass density."),
            ("B", "Content B is also reasonably long enough."),
            ("C", "Content C has enough characters for scoring."),
        ],
    )


def _make_poor_meta(path: Path | None = None) -> PresentationMeta:
    """A low-scoring extraction result."""
    return PresentationMeta(
        title="Presentation",  # generic
        file_path=path or Path("/tmp/deck.pptx"),
        slide_count=6,
        topics=[],
        slide_summaries=[],
    )


class TestBuildExtractorChain:
    def test_pptx_includes_pptx_and_ooxml(self, tmp_path):
        chain = _build_extractor_chain(tmp_path / "deck.pptx")
        names = [name for name, _ in chain]
        assert "PptxExtractor" in names
        assert "OoxmlExtractor" in names

    @patch("slidesmd.extractor._get_libreoffice_extractor", return_value=None)
    def test_ppt_without_libreoffice_raises(self, mock_lo):
        with pytest.raises(RuntimeError, match="No extractor available for .ppt"):
            _build_extractor_chain(Path("/tmp/legacy.ppt"))

    @patch("slidesmd.extractor._get_libreoffice_extractor")
    def test_ppt_with_libreoffice(self, mock_lo):
        mock_lo.return_value = MagicMock()
        chain = _build_extractor_chain(Path("/tmp/legacy.ppt"))
        names = [name for name, _ in chain]
        assert "LibreOfficeExtractor" in names
        assert "PptxExtractor" not in names


class TestCascadeLogic:
    @patch("slidesmd.extractor._build_extractor_chain")
    def test_returns_first_high_quality_result(self, mock_chain):
        good = _make_meta(title="Q3 Review")
        ext1 = MagicMock()
        ext1.extract.return_value = good

        ext2 = MagicMock()
        mock_chain.return_value = [("Ext1", ext1), ("Ext2", ext2)]

        result = _cascading_extract(Path("/tmp/deck.pptx"))
        assert result.title == "Q3 Review"
        ext2.extract.assert_not_called()  # never reached

    @patch("slidesmd.extractor._build_extractor_chain")
    def test_falls_through_to_second_extractor(self, mock_chain):
        poor = _make_poor_meta()
        good = _make_meta(title="Recovered Content")

        ext1 = MagicMock()
        ext1.extract.return_value = poor

        ext2 = MagicMock()
        ext2.extract.return_value = good

        mock_chain.return_value = [("Ext1", ext1), ("Ext2", ext2)]

        result = _cascading_extract(Path("/tmp/deck.pptx"))
        assert result.title == "Recovered Content"

    @patch("slidesmd.extractor._build_extractor_chain")
    def test_returns_best_when_none_meet_threshold(self, mock_chain):
        # Both are poor: generic title, 10 slides but very few summaries
        poor1 = PresentationMeta(
            title="Presentation", file_path=Path("/tmp/deck.pptx"),
            slide_count=10, topics=["A"], slide_summaries=[("A", "Short")],
        )
        poor2 = PresentationMeta(
            title="Presentation", file_path=Path("/tmp/deck.pptx"),
            slide_count=10, topics=[], slide_summaries=[],
        )

        ext1 = MagicMock()
        ext1.extract.return_value = poor1

        ext2 = MagicMock()
        ext2.extract.return_value = poor2

        mock_chain.return_value = [("Ext1", ext1), ("Ext2", ext2)]

        result = _cascading_extract(Path("/tmp/deck.pptx"))
        # poor1 has more content than poor2, so it should be returned
        assert result is poor1

    @patch("slidesmd.extractor._build_extractor_chain")
    def test_skips_failed_extractor(self, mock_chain):
        good = _make_meta(title="Fallback Result")

        ext1 = MagicMock()
        ext1.extract.side_effect = RuntimeError("python-pptx crashed")

        ext2 = MagicMock()
        ext2.extract.return_value = good

        mock_chain.return_value = [("Ext1", ext1), ("Ext2", ext2)]

        result = _cascading_extract(Path("/tmp/deck.pptx"))
        assert result.title == "Fallback Result"

    @patch("slidesmd.extractor._build_extractor_chain")
    def test_raises_when_all_fail(self, mock_chain):
        ext1 = MagicMock()
        ext1.extract.side_effect = RuntimeError("fail 1")

        ext2 = MagicMock()
        ext2.extract.side_effect = RuntimeError("fail 2")

        mock_chain.return_value = [("Ext1", ext1), ("Ext2", ext2)]

        with pytest.raises(RuntimeError, match="All extractors failed"):
            _cascading_extract(Path("/tmp/deck.pptx"))

    @patch("slidesmd.extractor._build_extractor_chain")
    def test_custom_threshold_via_env(self, mock_chain, monkeypatch):
        """Setting SLIDESMD_QUALITY_THRESHOLD=0.99 forces fallthrough."""
        monkeypatch.setenv("SLIDESMD_QUALITY_THRESHOLD", "0.99")

        good = _make_meta(title="Q3 Review")  # scores ~0.8
        better = _make_meta(title="Full Recovery", summaries=[
            ("A", "x" * 100), ("B", "x" * 100), ("C", "x" * 100),
        ])

        ext1 = MagicMock()
        ext1.extract.return_value = good

        ext2 = MagicMock()
        ext2.extract.return_value = better

        mock_chain.return_value = [("Ext1", ext1), ("Ext2", ext2)]

        _cascading_extract(Path("/tmp/deck.pptx"))
        # Even 'good' doesn't meet 0.99, so cascade tries ext2
        ext2.extract.assert_called_once()


class TestExtractPublicAPI:
    def test_explicit_extractor_bypasses_cascade(self):
        custom = MagicMock()
        custom.extract.return_value = _make_meta(title="Custom")

        result = extract(Path("/tmp/deck.pptx"), extractor=custom)
        assert result.title == "Custom"

    def test_default_uses_cascade(self, tmp_path):
        """Without explicit extractor, cascade is used (via real python-pptx)."""
        from pptx import Presentation

        pptx = tmp_path / "test.pptx"
        prs = Presentation()
        prs.core_properties.title = "Cascade Test"
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Slide 1"
        slide.placeholders[1].text = "Some content for the slide body."
        prs.save(str(pptx))

        result = extract(pptx)
        assert result.title == "Cascade Test"
        assert result.slide_count == 1
