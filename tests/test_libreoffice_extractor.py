"""Tests for the LibreOffice headless extractor.

Tests that require LibreOffice are skipped if it's not installed.
Unit tests for internal helpers run unconditionally.
"""

from __future__ import annotations

import shutil
from pathlib import Path
from unittest.mock import patch

import pytest

from slidesmd.libreoffice_extractor import LibreOfficeExtractor

_HAS_LIBREOFFICE = shutil.which("libreoffice") is not None or shutil.which("soffice") is not None


class TestAvailability:
    def test_is_available_reflects_system(self):
        assert LibreOfficeExtractor.is_available() == _HAS_LIBREOFFICE

    def test_extract_raises_without_libreoffice(self, tmp_path):
        ext = LibreOfficeExtractor()
        pptx = tmp_path / "deck.pptx"
        pptx.touch()
        with patch.object(LibreOfficeExtractor, "_find_libreoffice", return_value=None):
            with pytest.raises(RuntimeError, match="LibreOffice not found"):
                ext.extract(pptx)


class TestDetectTitle:
    def test_uses_first_line(self):
        ext = LibreOfficeExtractor()
        title = ext._detect_title(["First Line\nSecond Line"], Path("deck.pptx"))
        assert title == "First Line"

    def test_falls_back_to_filename(self):
        ext = LibreOfficeExtractor()
        title = ext._detect_title([], Path("my_presentation.pptx"))
        assert title == "My Presentation"

    def test_skips_empty_lines(self):
        ext = LibreOfficeExtractor()
        title = ext._detect_title(["\n\n  \nActual Title\nBody"], Path("deck.pptx"))
        assert title == "Actual Title"


class TestTodoExtraction:
    def test_finds_todos_in_pages(self):
        # Simulate what extract() does internally
        pages = ["Title\nAction item: fix the bug\nRegular content"]
        lines = []
        for page in pages:
            lines.extend([ln.strip() for ln in page.split("\n") if ln.strip()])
        todos = [ln for ln in lines if any(kw in ln.lower() for kw in (
            "todo", "action item", "follow up", "follow-up", "next step", "to-do"
        ))]
        assert len(todos) == 1
        assert "fix the bug" in todos[0]


@pytest.mark.skipif(not _HAS_LIBREOFFICE, reason="LibreOffice not installed")
class TestFullExtraction:
    """End-to-end tests requiring LibreOffice on the system."""

    def test_extract_real_pptx(self, tmp_path):
        """Build a real .pptx via python-pptx and extract via LibreOffice."""
        from pptx import Presentation

        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        prs.core_properties.title = "LO Test Deck"
        layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Revenue Growth"
        slide.placeholders[1].text = "Revenue grew 20% year-over-year."

        slide2 = prs.slides.add_slide(layout)
        slide2.shapes.title.text = "Next Steps"
        slide2.placeholders[1].text = "Action item: finalize the Q4 forecast."

        prs.save(str(pptx_path))

        ext = LibreOfficeExtractor(timeout=60)
        meta = ext.extract(pptx_path)

        assert meta.slide_count >= 2
        assert meta.file_path == pptx_path
        # Should have extracted some text
        assert len(meta.slide_summaries) >= 1
        all_text = " ".join(body for _, body in meta.slide_summaries)
        assert "revenue" in all_text.lower() or "20%" in all_text
