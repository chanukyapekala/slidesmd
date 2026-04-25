"""End-to-end integration tests for the semantic search pipeline.

These tests use real .pptx files (generated via python-pptx), real
sentence-transformers embeddings, and real SQLite + numpy retrieval.
No mocking — the full pipeline is exercised.

NOTE: First run downloads the all-MiniLM-L6-v2 model (~80 MB).
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from slidesmd import embedder, vector_store
from slidesmd.extractor import extract


# ---------------------------------------------------------------------------
# Helpers — build real .pptx files programmatically
# ---------------------------------------------------------------------------


def _make_pptx(path: Path, title: str, slides: list[tuple[str, str]]) -> Path:
    """Create a .pptx file with the given title and (slide_title, body) pairs."""
    prs = Presentation()
    prs.core_properties.title = title

    for slide_title, body in slides:
        layout = prs.slide_layouts[1]  # title + content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = slide_title
        slide.placeholders[1].text = body

    prs.save(str(path))
    return path


@pytest.fixture()
def finance_deck(tmp_path: Path) -> Path:
    return _make_pptx(
        tmp_path / "finance.pptx",
        "Q3 Financial Review",
        [
            ("Revenue Performance", "Revenue grew 15% year-over-year to $42M. "
             "Subscription revenue increased 22%. Enterprise deals up 30%."),
            ("Profit Margins", "Gross margin improved to 72%. EBITDA margin reached 22%. "
             "Operating expenses reduced by 5% through automation."),
            ("Cash Flow", "Free cash flow of $8.5M. Cash reserves at $120M. "
             "No debt on balance sheet."),
            ("Guidance", "Full year revenue target of $180M. Expecting 18-20% growth. "
             "Planning to increase R&D investment by 15%."),
        ],
    )


@pytest.fixture()
def product_deck(tmp_path: Path) -> Path:
    return _make_pptx(
        tmp_path / "product.pptx",
        "Product Strategy 2024",
        [
            ("Mobile App Launch", "Launching iOS and Android apps in Q4. "
             "Feature parity with web by end of year. Beta testing starts in September."),
            ("AI Features Roadmap", "Adding natural language search powered by embeddings. "
             "Smart summarization of documents. Auto-tagging and classification."),
            ("Competitive Landscape", "Main competitors: Acme Corp and Beta Inc. "
             "Our differentiation: better developer experience and lower latency."),
        ],
    )


@pytest.fixture()
def hr_deck(tmp_path: Path) -> Path:
    return _make_pptx(
        tmp_path / "hr_update.pptx",
        "HR & People Update",
        [
            ("Hiring Update", "Hired 12 new engineers this quarter. "
             "Attrition rate dropped to 8%. Engineering headcount now at 85."),
            ("Employee Satisfaction", "eNPS score improved to 72. "
             "Top concerns: career growth paths and remote work flexibility."),
            ("Diversity Metrics", "Women in engineering increased to 35%. "
             "Launched mentorship program for underrepresented groups."),
        ],
    )


# ---------------------------------------------------------------------------
# End-to-end tests
# ---------------------------------------------------------------------------


class TestFullPipeline:
    """Test the complete: extract → chunk → embed → store → search cycle."""

    def test_index_and_search_single_deck(self, finance_deck, tmp_path):
        """Index one deck and verify semantic search finds relevant slides."""
        indexed, skipped = embedder.index_folder(tmp_path)
        assert indexed == 1
        assert skipped == 0

        # search for something semantically related to revenue
        results = embedder.search(tmp_path, "How much money did the company make?", top_k=2)
        assert len(results) > 0

        # the top result should be from the revenue or profit slide
        top = results[0]
        assert top.chunk.file_path == str(finance_deck)
        assert top.chunk.slide_title in ("Revenue Performance", "Profit Margins")
        assert top.distance < 0.7  # cosine distance; <1.0 means positively correlated

    def test_index_multiple_decks(self, finance_deck, product_deck, hr_deck, tmp_path):
        """Index three decks and verify cross-deck semantic search."""
        indexed, skipped = embedder.index_folder(tmp_path)
        assert indexed == 3
        assert skipped == 0

        # verify total chunks in DB
        conn = vector_store.open_store(tmp_path)
        count = conn.execute("SELECT count(*) FROM chunks").fetchone()[0]
        assert count == 10  # 4 + 3 + 3 slides
        conn.close()

    def test_search_finds_correct_deck(self, finance_deck, product_deck, hr_deck, tmp_path):
        """Semantic search should rank the right deck highest."""
        embedder.index_folder(tmp_path)

        # query about hiring should surface HR deck
        results = embedder.search(tmp_path, "How many people did we hire?", top_k=3)
        assert len(results) >= 1
        top = results[0]
        assert Path(top.chunk.file_path).name == "hr_update.pptx"
        assert "Hiring" in top.chunk.slide_title or "engineer" in top.chunk.chunk_text.lower()

        # query about mobile app should surface product deck
        results = embedder.search(tmp_path, "When is the mobile app launching?", top_k=3)
        assert len(results) >= 1
        top = results[0]
        assert Path(top.chunk.file_path).name == "product.pptx"

        # query about cash should surface finance deck
        results = embedder.search(tmp_path, "What are our cash reserves?", top_k=3)
        assert len(results) >= 1
        top = results[0]
        assert Path(top.chunk.file_path).name == "finance.pptx"

    def test_ranking_quality(self, finance_deck, product_deck, hr_deck, tmp_path):
        """Top result should have a better (lower) distance than the last."""
        embedder.index_folder(tmp_path)

        results = embedder.search(tmp_path, "artificial intelligence and machine learning features", top_k=5)
        assert len(results) >= 2
        # results should be sorted by distance (ascending)
        distances = [r.distance for r in results]
        assert distances == sorted(distances)
        # best match should be significantly better than worst
        assert results[0].distance < results[-1].distance


class TestIncrementalIndex:
    """Test that re-indexing is incremental (only changed files re-embedded)."""

    def test_skip_unchanged(self, finance_deck, tmp_path):
        """Second index of the same file should skip it."""
        indexed1, _ = embedder.index_folder(tmp_path)
        assert indexed1 == 1

        indexed2, skipped2 = embedder.index_folder(tmp_path)
        assert indexed2 == 0
        assert skipped2 == 1

    def test_reindex_after_modification(self, tmp_path):
        """If a file is modified, it should be re-embedded."""
        import time

        deck = _make_pptx(
            tmp_path / "evolving.pptx",
            "V1",
            [("Slide 1", "Original content about cats and dogs.")],
        )
        indexed1, _ = embedder.index_folder(tmp_path)
        assert indexed1 == 1

        # wait briefly and overwrite with new content
        time.sleep(0.05)
        _make_pptx(
            tmp_path / "evolving.pptx",
            "V2",
            [("Slide 1", "Updated content about rockets and space exploration.")],
        )

        indexed2, skipped2 = embedder.index_folder(tmp_path)
        assert indexed2 == 1
        assert skipped2 == 0

        # search should find the new content
        results = embedder.search(tmp_path, "space exploration", top_k=1)
        assert len(results) == 1
        assert "rocket" in results[0].chunk.chunk_text.lower() or "space" in results[0].chunk.chunk_text.lower()

    def test_adding_new_deck(self, finance_deck, tmp_path):
        """Adding a new deck should only index the new one."""
        embedder.index_folder(tmp_path)

        _make_pptx(
            tmp_path / "new_deck.pptx",
            "Brand New",
            [("Fresh Content", "This is completely new material about quantum computing.")],
        )

        indexed, skipped = embedder.index_folder(tmp_path)
        assert indexed == 1  # only the new file
        assert skipped == 1  # finance deck unchanged


class TestVectorStoreIntegrity:
    """Verify the SQLite store contents directly after indexing."""

    def test_chunk_metadata_stored_correctly(self, finance_deck, tmp_path):
        """Each chunk should have correct metadata in the DB."""
        embedder.index_folder(tmp_path)

        conn = vector_store.open_store(tmp_path)
        rows = conn.execute(
            "SELECT file_path, slide_index, slide_title, chunk_text FROM chunks ORDER BY slide_index"
        ).fetchall()
        conn.close()

        assert len(rows) == 4  # 4 slides in finance deck
        assert rows[0][2] == "Revenue Performance"
        assert rows[1][2] == "Profit Margins"
        assert rows[2][2] == "Cash Flow"
        assert rows[3][2] == "Guidance"

        # chunk text should contain the presentation title prefix
        for row in rows:
            assert "Q3 Financial Review" in row[3]

    def test_embedding_blob_correct_size(self, finance_deck, tmp_path):
        """Stored embedding blobs should be 384 * 4 bytes."""
        embedder.index_folder(tmp_path)

        conn = vector_store.open_store(tmp_path)
        rows = conn.execute("SELECT embedding FROM chunks").fetchall()
        conn.close()

        for (blob,) in rows:
            assert len(blob) == vector_store.EMBEDDING_DIM * 4  # 384 floats * 4 bytes

    def test_delete_removes_all_chunks(self, finance_deck, tmp_path):
        """Deleting a file should remove all its chunks."""
        embedder.index_folder(tmp_path)

        conn = vector_store.open_store(tmp_path)
        vector_store.delete_file_chunks(conn, str(finance_deck))
        count = conn.execute("SELECT count(*) FROM chunks").fetchone()[0]
        conn.close()
        assert count == 0
