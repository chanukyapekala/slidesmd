"""Extract metadata and to-dos from .pptx files."""

from __future__ import annotations

import logging
import os
from dataclasses import dataclass, field
from pathlib import Path
from typing import Protocol, runtime_checkable

from slidesmd.image_parser import ImageResult, extract_images_from_slide, parse_image

logger = logging.getLogger(__name__)


@dataclass
class PresentationMeta:
    title: str
    file_path: Path
    slide_count: int
    topics: list[str] = field(default_factory=list)
    todos: list[str] = field(default_factory=list)
    slide_summaries: list[tuple[str, str]] = field(default_factory=list)  # (title, body)
    image_results: list[tuple[str, ImageResult]] = field(default_factory=list)  # (slide_title, result)


@runtime_checkable
class SlideExtractor(Protocol):
    """Protocol for pluggable PPTX extraction backends."""

    def extract(self, path: Path) -> PresentationMeta:
        """Extract structured metadata from a presentation file."""
        ...


class PptxExtractor:
    """Default extractor backed by python-pptx."""

    def extract(self, path: Path) -> PresentationMeta:
        from pptx import Presentation  # deferred to avoid import at module load time; custom extractors need not use pptx at all

        prs = Presentation(path)
        return PresentationMeta(
            title=_extract_title(prs, path),
            file_path=path,
            slide_count=len(prs.slides),
            topics=_extract_topics(prs),
            todos=_extract_todos(prs),
            slide_summaries=_extract_slide_summaries(prs),
            image_results=_extract_images(prs),
        )


_default_extractor = PptxExtractor()

_QUALITY_THRESHOLD = 0.4


def extract(pptx_path: Path, extractor: SlideExtractor | None = None) -> PresentationMeta:
    """Extract metadata from a presentation file.

    Pass a custom *extractor* to bypass the cascade and use it directly.
    When no extractor is specified, a scored cascade is used:
    PptxExtractor → OoxmlExtractor → LibreOfficeExtractor.
    """
    if extractor is not None:
        return extractor.extract(pptx_path)
    return _cascading_extract(pptx_path)


# ---------------------------------------------------------------------------
# Cascade logic
# ---------------------------------------------------------------------------

_ooxml_instance: object | None = None
_lo_instance: object | None = None


def _get_ooxml_extractor() -> SlideExtractor:
    global _ooxml_instance
    if _ooxml_instance is None:
        from slidesmd.ooxml_extractor import OoxmlExtractor
        _ooxml_instance = OoxmlExtractor()
    return _ooxml_instance  # type: ignore[return-value]


def _get_libreoffice_extractor() -> SlideExtractor | None:
    global _lo_instance
    if _lo_instance is None:
        from slidesmd.libreoffice_extractor import LibreOfficeExtractor
        if LibreOfficeExtractor.is_available():
            _lo_instance = LibreOfficeExtractor()
        else:
            return None
    return _lo_instance  # type: ignore[return-value]


def _build_extractor_chain(path: Path) -> list[tuple[str, SlideExtractor]]:
    """Select extractors appropriate for the file type."""
    suffix = path.suffix.lower()
    chain: list[tuple[str, SlideExtractor]] = []

    if suffix == ".pptx":
        chain.append(("PptxExtractor", _default_extractor))
        chain.append(("OoxmlExtractor", _get_ooxml_extractor()))

    lo = _get_libreoffice_extractor()
    if lo is not None:
        chain.append(("LibreOfficeExtractor", lo))

    if not chain:
        if suffix in (".ppt", ".odp"):
            raise RuntimeError(
                f"No extractor available for {suffix} files (install LibreOffice)"
            )
        # Unknown suffix, try PptxExtractor as best effort
        chain.append(("PptxExtractor", _default_extractor))

    return chain


def _cascading_extract(path: Path) -> PresentationMeta:
    """Try extractors in order, scoring each result."""
    from slidesmd.scorer import score_extraction

    threshold = float(os.environ.get("SLIDESMD_QUALITY_THRESHOLD", str(_QUALITY_THRESHOLD)))
    chain = _build_extractor_chain(path)

    best_result: PresentationMeta | None = None
    best_score: float = -1.0

    for name, ext in chain:
        try:
            result = ext.extract(path)
            score = score_extraction(result)
            logger.debug("Extractor %s scored %.2f for %s", name, score, path.name)

            if score >= threshold:
                return result

            if score > best_score:
                best_score = score
                best_result = result
        except Exception as e:
            logger.debug("Extractor %s failed for %s: %s", name, path.name, e)
            continue

    if best_result is not None:
        return best_result

    raise RuntimeError(f"All extractors failed for {path}")


# ---------------------------------------------------------------------------
# Internal helpers (python-pptx specific — used only by PptxExtractor)
# ---------------------------------------------------------------------------

def _placeholder_idx(shape: object) -> int | None:
    """Return placeholder index or None if shape is not a placeholder."""
    try:
        fmt = shape.placeholder_format  # type: ignore[attr-defined]
        return fmt.idx if fmt is not None else None
    except Exception:
        return None


_GENERIC_TITLES = {"powerpoint presentation", "presentation", "untitled"}


def _first_slide_title(prs: object) -> str:
    try:
        slides = prs.slides  # type: ignore[attr-defined]
    except AttributeError:
        return ""
    if slides:
        for shape in slides[0].shapes:
            if _placeholder_idx(shape) == 0 and shape.has_text_frame:
                return shape.text_frame.text.strip()
    return ""


def _extract_title(prs: object, fallback: Path) -> str:
    """Use core properties title unless it's generic, then prefer first slide title."""
    core_title = (prs.core_properties.title or "").strip()  # type: ignore[attr-defined]
    if core_title and core_title.lower() not in _GENERIC_TITLES:
        return core_title

    slide_title = _first_slide_title(prs)
    if slide_title:
        return slide_title

    if core_title:
        return core_title

    return fallback.stem.replace("-", " ").replace("_", " ").title()


def _extract_topics(prs: object) -> list[str]:
    """Extract slide titles as topic list."""
    topics: list[str] = []
    for slide in prs.slides:  # type: ignore[attr-defined]
        for shape in slide.shapes:
            if _placeholder_idx(shape) == 0 and shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    topics.append(text)
    return topics


def _extract_slide_summaries(prs: object) -> list[tuple[str, str]]:
    """Extract (slide_title, body_text) for each slide."""
    summaries = []
    for slide in prs.slides:  # type: ignore[attr-defined]
        slide_title = ""
        body_parts: list[str] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            idx = _placeholder_idx(shape)
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if idx == 0:
                slide_title = text
            else:
                body_parts.append(text)

        body = " · ".join(body_parts)
        if slide_title or body:
            summaries.append((slide_title, body))

    return summaries


def _extract_images(prs: object) -> list[tuple[str, ImageResult]]:
    """Extract and parse images from all slides."""
    results = []
    for slide in prs.slides:  # type: ignore[attr-defined]
        slide_title = ""
        for shape in slide.shapes:
            if _placeholder_idx(shape) == 0 and shape.has_text_frame:
                slide_title = shape.text_frame.text.strip()
                break

        for image in extract_images_from_slide(slide):
            result = parse_image(image, slide_title)
            if result.method != "skipped":
                results.append((slide_title, result))

    return results


def _extract_todos(prs: object) -> list[str]:
    """Extract lines containing TODO, Action, or follow-up keywords."""
    keywords = ("todo", "action item", "follow up", "follow-up", "next step", "to-do")
    todos: list[str] = []

    for slide in prs.slides:  # type: ignore[attr-defined]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if any(kw in text.lower() for kw in keywords):
                    todos.append(text)

    return todos
